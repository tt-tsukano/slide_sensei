import os
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation
from langchain.schema import Document

def load_presentation(file_path):
    """
    PDFまたはPPTXファイルを読み込み、ページ/スライド単位でDocumentオブジェクトを返す
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        return documents

    elif ext == ".pptx":
        prs = Presentation(file_path)
        documents = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
            slide_text = "\n".join(texts).strip()
            if slide_text:
                documents.append(Document(
                    page_content=slide_text,
                    metadata={"slide_number": i+1}
                ))
        return documents

    else:
        raise ValueError("対応していないファイル形式です（PDFまたはPPTXのみ対応）")

if __name__ == "__main__":
    test_file = "sample.pptx"  # または sample.pdf
    if not os.path.exists(test_file):
        print(f"ファイルが見つかりません: {test_file}")
    else:
        docs = load_presentation(test_file)
        for doc in docs:
            print(f"--- Slide/Page {doc.metadata.get('slide_number', '?')} ---")
            print(doc.page_content[:200])
            print()